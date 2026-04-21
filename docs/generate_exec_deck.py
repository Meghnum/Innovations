"""Generate a 3-slide executive deck for the Claims Assistant.

Audience: directors / seniors. Story arc:
  1. Ask    — Claims Assistant (natural-language analytics)
  2. Decide — Fast Track Triage + Rule Optimizer (learning loop)
  3. Prove  — Shadow-Mode KPIs (safety before rollout)

Output: docs/Claims_Assistant_Exec_Deck.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Palette (dark, modern, single accent per slide) ────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)   # near-black
SURFACE   = RGBColor(0x16, 0x1B, 0x22)   # tile surface
BORDER    = RGBColor(0x30, 0x36, 0x3D)
TEXT      = RGBColor(0xE6, 0xED, 0xF3)
SUBTLE    = RGBColor(0x8B, 0x94, 0x9E)
GREEN     = RGBColor(0x3F, 0xB9, 0x50)   # slide 1 accent — Ask
AMBER     = RGBColor(0xD2, 0x99, 0x22)   # slide 2 accent — Decide
BLUE      = RGBColor(0x58, 0xA6, 0xFF)   # slide 3 accent — Prove
RED       = RGBColor(0xF8, 0x51, 0x49)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)   # 16:9 widescreen


def _add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False


def _text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _tile(slide, x, y, w, h, accent):
    t = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    t.adjustments[0] = 0.06
    t.fill.solid(); t.fill.fore_color.rgb = SURFACE
    t.line.color.rgb = BORDER
    t.line.width = Pt(0.75)
    t.shadow.inherit = False
    # Accent bar on the left edge
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    return t


def _header(slide, kicker, title, accent):
    # Small colored kicker strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.55),
                                   Inches(0.35), Inches(0.07))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    _text(slide, Inches(1.05), Inches(0.45), Inches(6), Inches(0.3),
          kicker.upper(), size=11, bold=True, color=accent)
    _text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
          title, size=34, bold=True, color=TEXT)


def _footer(slide, note, page, total):
    _text(slide, Inches(0.6), Inches(7.0), Inches(10), Inches(0.35),
          note, size=10, color=SUBTLE)
    _text(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.35),
          f"{page} / {total}", size=10, color=SUBTLE, align=PP_ALIGN.RIGHT)


def _bullet_block(slide, x, y, w, h, icon, title, body, accent):
    _tile(slide, x, y, w, h, accent)
    _text(slide, x + Inches(0.3), y + Inches(0.2),
          Inches(0.7), Inches(0.5), icon, size=28, color=accent)
    _text(slide, x + Inches(0.3), y + Inches(0.75),
          w - Inches(0.5), Inches(0.45), title,
          size=16, bold=True, color=TEXT)
    _text(slide, x + Inches(0.3), y + Inches(1.22),
          w - Inches(0.5), h - Inches(1.3), body,
          size=12, color=SUBTLE)


# ── Build ──────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    # ──────────────────────────────────────────────────────────────────────
    # Slide 1 — ASK  (Claims Assistant)
    # ──────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank); _add_bg(s1)
    _header(s1, "01 · Ask", "Claims Assistant — talk to 50,000 claims in English",
            GREEN)

    # Left: the problem → solution bubble
    _text(s1, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.5),
          "The problem", size=13, bold=True, color=GREEN)
    _text(s1, Inches(0.6), Inches(2.15), Inches(6.0), Inches(1.3),
          "Adjusters and managers wait days for analysts to pull numbers.\n"
          "Spreadsheets age the moment they're exported.",
          size=15, color=TEXT)

    _text(s1, Inches(0.6), Inches(3.7), Inches(6.0), Inches(0.5),
          "What we built", size=13, bold=True, color=GREEN)
    _text(s1, Inches(0.6), Inches(4.05), Inches(6.0), Inches(2.4),
          "A chat assistant that reads the live claims dataset.\n\n"
          "•  Plain-English questions → instant answers\n"
          "•  Hybrid engine: rules for common asks, LLM for complex math\n"
          "•  Vector search across every loss description\n"
          "•  Runs locally — no claims data leaves our network",
          size=14, color=TEXT)

    # Right: chat-mock tile
    box = _tile(s1, Inches(7.4), Inches(1.8), Inches(5.3), Inches(4.8), GREEN)
    _text(s1, Inches(7.65), Inches(1.95), Inches(5.0), Inches(0.4),
          "EXAMPLE", size=10, bold=True, color=SUBTLE)
    _text(s1, Inches(7.65), Inches(2.3), Inches(5.0), Inches(0.5),
          "> Top 5 LOBs by incurred, excluding Casualty",
          size=13, bold=True, color=TEXT, font="Consolas")
    _text(s1, Inches(7.65), Inches(2.95), Inches(5.0), Inches(3.5),
          "Property         $412.7M\n"
          "Marine           $298.1M\n"
          "Cyber            $187.4M\n"
          "A&H              $154.8M\n"
          "Auto             $122.3M\n\n"
          "Answer generated in 1.8s — cites 5 source claims.",
          size=12, color=TEXT, font="Consolas")

    _footer(s1, "Natural-language analytics over the full book, on-prem.", 1, 3)

    # ──────────────────────────────────────────────────────────────────────
    # Slide 2 — DECIDE  (Fast Track Triage + Rule Optimizer)
    # ──────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank); _add_bg(s2)
    _header(s2, "02 · Decide", "Fast-Track Triage + Rule Optimizer — a system that learns",
            AMBER)

    # Three stacked tiles across the slide
    y = Inches(1.9); h = Inches(1.55); gap = Inches(0.2)

    _bullet_block(
        s2, Inches(0.6), y, Inches(12.1), h, "▶",
        "Fast-Track Triage — every new claim gets a verdict in seconds",
        "Hard rules (reserve, LOB, injury severity) + vector precedent memory "
        "recommend FAST TRACK or MANUAL REVIEW with a plain-English reason.",
        AMBER,
    )

    _bullet_block(
        s2, Inches(0.6), y + h + gap, Inches(12.1), h, "↺",
        "Human-in-the-loop — adjusters approve or override the AI",
        "One click writes the decision back into the precedent index. Next "
        "similar claim sees the new pattern instantly — no retraining cycle.",
        AMBER,
    )

    _bullet_block(
        s2, Inches(0.6), y + 2 * (h + gap), Inches(12.1), h, "⚙",
        "Rule Optimizer — flags which rules to tune, before the board does",
        "Watches the disagreement log, surfaces candidate threshold changes "
        "(e.g. '$7.5k → $9k reserve cap'), manager approves with one click.",
        AMBER,
    )

    _footer(s2, "Decide → feedback → learn. The loop closes on every claim.", 2, 3)

    # ──────────────────────────────────────────────────────────────────────
    # Slide 3 — PROVE  (Shadow-Mode KPIs)
    # ──────────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank); _add_bg(s3)
    _header(s3, "03 · Prove", "Shadow Mode — measured safety before we go live",
            BLUE)

    _text(s3, Inches(0.6), Inches(1.75), Inches(12), Inches(0.6),
          "The AI runs alongside adjusters — no automation yet. "
          "Every decision pair is logged and scored against three targets:",
          size=14, color=TEXT)

    # Three KPI tiles side-by-side
    kpi_y = Inches(2.8); kpi_h = Inches(3.3); kpi_w = Inches(3.95); x0 = Inches(0.6); gap = Inches(0.25)

    def kpi_tile(x, label, target, value, status_color, narrative):
        _tile(s3, x, kpi_y, kpi_w, kpi_h, status_color)
        _text(s3, x + Inches(0.3), kpi_y + Inches(0.25),
              kpi_w - Inches(0.5), Inches(0.45),
              label.upper(), size=11, bold=True, color=SUBTLE)
        _text(s3, x + Inches(0.3), kpi_y + Inches(0.7),
              kpi_w - Inches(0.5), Inches(1.3), value,
              size=48, bold=True, color=status_color)
        _text(s3, x + Inches(0.3), kpi_y + Inches(2.0),
              kpi_w - Inches(0.5), Inches(0.4),
              f"Target: {target}", size=12, bold=True, color=TEXT)
        _text(s3, x + Inches(0.3), kpi_y + Inches(2.4),
              kpi_w - Inches(0.5), Inches(0.8),
              narrative, size=11, color=SUBTLE)

    kpi_tile(x0,
             "Agreement Rate", "> 90%", ">90%", GREEN,
             "How often adjusters agree with the AI's recommendation.")
    kpi_tile(x0 + kpi_w + gap,
             "Leakage (False Positive)", "0%", "0%", RED,
             "AI said FAST TRACK but adjuster disagreed — paying something we shouldn't.")
    kpi_tile(x0 + 2 * (kpi_w + gap),
             "Friction (False Negative)", "< 10%", "<10%", AMBER,
             "AI over-cautious — flagged to MANUAL but adjuster would have fast-tracked.")

    _text(s3, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5),
          "Only when all three targets hold for a full 2-week window do we propose turning automation on.",
          size=13, bold=True, color=TEXT)

    _footer(s3, "We earn the right to automate. Evidence first, rollout second.", 3, 3)

    out = Path(__file__).parent / "Claims_Assistant_Exec_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
