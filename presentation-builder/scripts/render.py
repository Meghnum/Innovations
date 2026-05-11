from pathlib import Path
import matplotlib
matplotlib.use("Agg")  # non-interactive backend
import matplotlib.pyplot as plt
import seaborn as sns

import re
import polars as pl
from pptx.util import Pt
from pptx.dml.color import RGBColor
from difflib import get_close_matches

# ── chart constants ───────────────────────────────────────────────────────────
# Hardcoded fallback brand palette — replace with corporate colors via template extraction.
FALLBACK_COLORS = {
    "primary": "#1F4E79",
    "secondary": "#2E75B6",
    "accent": "#FFC000",
    "neutral": "#595959",
    "background": "#FFFFFF",
}


def get_brand_colors(template_path: str | None) -> dict:
    """Extract brand colors from corporate template's theme XML, or use fallback."""
    if template_path is None or not Path(template_path).exists():
        return dict(FALLBACK_COLORS)
    try:
        import zipfile
        with zipfile.ZipFile(template_path) as z:
            theme_files = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
            if not theme_files:
                return dict(FALLBACK_COLORS)
            theme_xml = z.read(theme_files[0]).decode("utf-8", errors="replace")
        # Extract sRGB hex codes from theme color scheme; first 6 are typical brand
        # Theme XML uses namespace; use regex to find srgbClr values reliably
        srgb_matches = re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', theme_xml)
        if len(srgb_matches) < 2:
            return dict(FALLBACK_COLORS)
        # Standard PPTX theme order: bg1, tx1, bg2, tx2, accent1..accent6
        # Use accent1 (index 4) as primary if available, else fallback
        colors = dict(FALLBACK_COLORS)
        # Try to map: tx1 (dark) → neutral, accent1 → primary, accent2 → secondary
        if len(srgb_matches) >= 5:
            colors["primary"] = "#" + srgb_matches[4].upper()
        if len(srgb_matches) >= 6:
            colors["secondary"] = "#" + srgb_matches[5].upper()
        if len(srgb_matches) >= 7:
            colors["accent"] = "#" + srgb_matches[6].upper()
        return colors
    except Exception:
        return dict(FALLBACK_COLORS)


def render_chart(chart_data: dict, out_path: str, title: str = "", template_path: str | None = None) -> dict:
    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])
    chart_type = chart_data.get("chart_type", "bar")
    if not labels or not values:
        return {"error": "no data to render", "png_path": None}

    colors = get_brand_colors(template_path)
    sns.set_style("whitegrid")
    fig, ax = plt.subplots(figsize=(10, 6), dpi=150)

    if chart_type == "line":
        ax.plot(labels, values, color=colors["primary"], linewidth=2.5, marker="o", markersize=8)
    elif chart_type == "bar":
        ax.bar(labels, values, color=colors["primary"], edgecolor=colors["neutral"])
    elif chart_type == "scatter":
        ax.scatter(labels, values, color=colors["primary"], s=80)
    elif chart_type == "stacked_bar":
        ax.bar(labels, values, color=colors["primary"])
    elif chart_type == "histogram":
        ax.hist(values, bins=min(20, len(values)), color=colors["primary"], edgecolor=colors["neutral"])
    else:
        ax.bar(labels, values, color=colors["primary"])

    if title:
        ax.set_title(title, fontsize=14, fontweight="bold", color=colors["neutral"])
    ax.set_xlabel(chart_data.get("x_axis", ""), color=colors["neutral"])
    ax.set_ylabel(chart_data.get("y_axis", ""), color=colors["neutral"])
    plt.xticks(rotation=30, ha="right")
    plt.tight_layout()

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path, dpi=300, bbox_inches="tight")
    plt.close(fig)
    return {"png_path": out_path, "width_px": 3000, "height_px": 1800}


# ── tables ────────────────────────────────────────────────────────────────────

MAX_ROWS = 10
MAX_COLS = 5
HEADER_FILL = RGBColor(0x1F, 0x4E, 0x79)
HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ALT_ROW_FILL = RGBColor(0xF2, 0xF2, 0xF2)


def add_native_table(slide, df: pl.DataFrame, left, top, width, height) -> dict:
    if df.height > MAX_ROWS or df.width > MAX_COLS:
        return {
            "error": f"slice exceeds 10x5 rule (rows={df.height}, cols={df.width})",
        }
    rows = df.height + 1
    cols = df.width
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Header
    for c, name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = HEADER_FONT_COLOR
                run.font.size = Pt(11)

    # Body
    data = df.to_dicts()
    for r, row in enumerate(data, start=1):
        for c, name in enumerate(df.columns):
            cell = table.cell(r, c)
            cell.text = str(row[name]) if row[name] is not None else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW_FILL

    return {
        "rows": rows,
        "cols": cols,
        "shape_id": table_shape.shape_id,
    }


# ── narrative ─────────────────────────────────────────────────────────────────

# Only validate numbers that look like financial/statistical claims:
# - $1,234.56 / $1234 / $1,234
# - 12.5% / -7%
# - 1,234,567 (with comma separators)
# Skip bare integers (years, ordinals, counts).
NUMBER_RX = re.compile(
    r"-?\$\d{1,3}(?:,\d{3})*(?:\.\d+)?"  # $1,234.56 or $1234
    r"|-?\d+(?:\.\d+)?%"                  # 12.5% or -7%
    r"|-?\d{1,3}(?:,\d{3})+(?:\.\d+)?"    # 1,234,567 (must have at least one comma)
)
TOLERANCE = 0.0001  # 0.01%
_NARRATIVE_SSN_RX = re.compile(r"\b\d{3}-\d{2}-\d{4}\b")
_NARRATIVE_EMAIL_RX = re.compile(r"\b[^@\s]+@[^@\s]+\.[^@\s]+\b")


def _normalize_number(s: str) -> float | None:
    s = s.replace("$", "").replace(",", "").replace("%", "")
    try:
        return float(s)
    except ValueError:
        return None


def _matches_kv(claim: float, kv_values: list) -> bool:
    for v in kv_values:
        if v == 0:
            if abs(claim) < 0.01:
                return True
            continue
        # also accept absolute value match (e.g. "Down 7%" referencing -7.0)
        if abs(claim - v) / abs(v) < TOLERANCE:
            return True
        if abs(v) != 0 and abs(abs(claim) - abs(v)) / abs(v) < TOLERANCE:
            return True
    return False


def validate_narrative(narrative: dict, kv: dict, pii_columns: list | None = None) -> dict:
    mismatches = []
    text = " ".join(str(narrative.get(k, "")) for k in ("observe", "analyze"))
    kv_values = [float(v) for v in kv.values() if isinstance(v, (int, float))]
    for raw in NUMBER_RX.findall(text):
        n = _normalize_number(raw)
        if n is None:
            continue
        if not _matches_kv(n, kv_values):
            mismatches.append(f"fabricated number: {raw}")

    # PII guards: scan all narrative tiers
    full_text = " ".join(str(narrative.get(k, "")) for k in ("observe", "analyze", "synthesize"))
    if pii_columns:
        for col in pii_columns:
            pattern = r"\b" + re.escape(col) + r"\b"
            if re.search(pattern, full_text, re.IGNORECASE):
                mismatches.append(f"PII column reference: {col}")
    if _NARRATIVE_SSN_RX.search(full_text):
        mismatches.append("PII: SSN pattern in narrative")
    if _NARRATIVE_EMAIL_RX.search(full_text):
        mismatches.append("PII: email pattern in narrative")

    return {"valid": len(mismatches) == 0, "mismatches": mismatches}


PROMPT_TEMPLATE = """You are an executive presentation analyst. Generate a 3-tier narrative for this slide.

SLIDE TITLE: {title}
AUDIENCE: {audience}

DATA FACTS (use ONLY these values; do not invent numbers):
{facts}

Produce exactly three statements following the Observe → Analyze → Synthesize chain:

1. OBSERVE: State the most important raw fact from the DATA FACTS (single sentence, include the literal number).
2. ANALYZE: State the comparative or trend context (single sentence, include the literal number for any delta or comparison).
3. SYNTHESIZE: State the business implication — the "So What?" (single sentence, no new numbers).

Constraints:
- Every numeric claim in OBSERVE and ANALYZE must appear verbatim in DATA FACTS above (rounding to nearest whole percent or dollar is allowed).
- Total ≤ 6 bullet points across the three tiers.
- Synthesize must be actionable for the AUDIENCE.
- Output as JSON: {{"observe": "...", "analyze": "...", "synthesize": "..."}}
"""


def build_prompt(kv: dict, slide_ctx: dict) -> str:
    facts_lines = [f"- {k}: {v}" for k, v in kv.items()]
    facts = "\n".join(facts_lines)
    return PROMPT_TEMPLATE.format(
        title=slide_ctx.get("title", "Untitled"),
        audience=slide_ctx.get("audience", "general"),
        facts=facts,
    )


# ── layouts ───────────────────────────────────────────────────────────────────

MAX_TABLE_ROWS = 10
MAX_TABLE_COLS = 5

LAYOUT_FOR_CONTENT = {
    "exec_summary": "Title and Content",
    "section": "Image + Text",
    "deep_dive": "Big Number",
    "table": "Table Layout",
    "title": "Title Slide",
}


def decide_render_mode(rows: int, cols: int) -> str:
    if rows <= MAX_TABLE_ROWS and cols <= MAX_TABLE_COLS:
        return "native_table"
    return "image"


def pick_layout(presentation, layout_name: str):
    """Find slide layout by name; fall back to closest match or first available."""
    available = {layout.name: layout for layout in presentation.slide_layouts}
    if layout_name in available:
        return available[layout_name]
    candidates = get_close_matches(layout_name, list(available.keys()), n=1, cutoff=0.3)
    if candidates:
        return available[candidates[0]]
    return presentation.slide_layouts[0]


def layout_for_content(content_type: str) -> str:
    return LAYOUT_FOR_CONTENT.get(content_type, "Title and Content")
