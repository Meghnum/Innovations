from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from scripts.render import pick_layout

DEFAULT_TEMPLATE = Path(__file__).parent.parent / "assets" / "default_template.pptx"


def _resolve_template(template_path: str | None) -> str:
    if template_path and Path(template_path).exists():
        return template_path
    if DEFAULT_TEMPLATE.exists():
        return str(DEFAULT_TEMPLATE)
    return None  # python-pptx falls back to its built-in default template


def _add_text(text_frame, text: str, font_size: int = 18):
    if not text_frame.paragraphs:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(font_size)


def _write_speaker_notes(slide, narrative: dict):
    notes_tf = slide.notes_slide.notes_text_frame
    parts = []
    if narrative.get("observe"):
        parts.append(f"Observe: {narrative['observe']}")
    if narrative.get("analyze"):
        parts.append(f"Analyze: {narrative['analyze']}")
    if narrative.get("synthesize"):
        parts.append(f"Synthesize: {narrative['synthesize']}")
    notes_tf.text = "\n\n".join(parts)


def _add_slide(prs, slide_data: dict):
    layout = pick_layout(prs, slide_data.get("layout", "Title and Content"))
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "")
    narrative = slide_data.get("narrative", {})

    body_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body_placeholder = shape
            break
    if body_placeholder and narrative.get("synthesize"):
        _add_text(body_placeholder.text_frame, narrative["synthesize"], font_size=18)
    elif narrative.get("synthesize"):
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        _add_text(tx.text_frame, narrative["synthesize"], font_size=18)

    if slide_data.get("chart_png"):
        slide.shapes.add_picture(
            slide_data["chart_png"],
            Inches(5), Inches(2),
            width=Inches(7),
        )

    if narrative:
        _write_speaker_notes(slide, narrative)


def _generate_exclusions_slide(prs, exclusions: list):
    if not exclusions:
        return
    layout = pick_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Data Integrity & Exclusions Report"
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    intro = "The following sections were excluded or modified to ensure accuracy and compliance:"
    lines = [intro] + [f"• {e['title']}: {e['reason']}" for e in exclusions]
    text = "\n".join(lines)
    if body:
        body.text_frame.text = text
    else:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tx.text_frame.text = text


def build_deck(outline: dict, template_path: str | None, out_path: str) -> dict:
    template = _resolve_template(template_path)
    prs = Presentation(template) if template else Presentation()
    exclusions = []
    for slide_data in outline.get("slides", []):
        if slide_data.get("status") == "excluded":
            exclusions.append({
                "title": slide_data.get("title", "(untitled)"),
                "reason": slide_data.get("reason", "no reason provided"),
            })
            continue
        _add_slide(prs, slide_data)
    _generate_exclusions_slide(prs, exclusions)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return {"pptx_path": out_path, "slide_count": len(prs.slides)}
