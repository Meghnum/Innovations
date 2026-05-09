import polars as pl
from pptx.util import Pt
from pptx.dml.color import RGBColor

MAX_ROWS = 10
MAX_COLS = 5
HEADER_FILL = RGBColor(0x1F, 0x4E, 0x79)
HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
ALT_ROW_FILL = RGBColor(0xF2, 0xF2, 0xF2)


def add_native_table(slide, df: pl.DataFrame, left, top, width, height) -> dict:
    if df.height > MAX_ROWS or df.width > MAX_COLS:
        return {
            "error": f"slice exceeds 10x5 rule (rows={df.height}, cols={df.width})",
        }
    rows = df.height + 1
    cols = df.width
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Header
    for c, name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = HEADER_FONT_COLOR
                run.font.size = Pt(11)

    # Body
    data = df.to_dicts()
    for r, row in enumerate(data, start=1):
        for c, name in enumerate(df.columns):
            cell = table.cell(r, c)
            cell.text = str(row[name]) if row[name] is not None else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW_FILL

    return {
        "rows": rows,
        "cols": cols,
        "shape_id": table_shape.shape_id,
    }
