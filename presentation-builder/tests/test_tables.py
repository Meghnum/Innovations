import polars as pl
from pptx import Presentation
from pptx.util import Inches
from scripts.render import add_native_table

def test_add_native_table_to_slide(tmp_path):
    df = pl.DataFrame({
        "Region": ["North", "South", "East"],
        "Revenue": [1000, 2000, 1500],
    })
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    descriptor = add_native_table(slide, df, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
    assert descriptor["rows"] == 4  # header + 3 data rows
    assert descriptor["cols"] == 2
    out = tmp_path / "out.pptx"
    prs.save(out)
    assert out.exists()

def test_table_rejects_oversized_slice():
    df = pl.DataFrame({f"c{i}": [0] for i in range(6)})  # 6 cols
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    descriptor = add_native_table(slide, df, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
    assert "error" in descriptor
    assert "exceeds" in descriptor["error"].lower()
