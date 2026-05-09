# presentation-builder/tests/test_e2e_messy.py
from pathlib import Path
from pptx import Presentation
from scripts.ingest import ingest
from scripts.profile import profile
from scripts.context import detect_context
from scripts.outline import build_outline
from scripts.build_pptx import build_deck


def test_high_null_file_produces_exclusions_slide(sample_messy, tmp_path):
    ing = ingest(str(sample_messy))
    df = ing["dataframe"]
    prof = profile(df)
    assert prof["null_pct"]["Cost"] > 15

    ctx = detect_context(prof)
    outline = build_outline(prof, ctx)
    margin_slide = [s for s in outline["slides"] if s.get("computation_id") == "gross_margin"]
    assert len(margin_slide) >= 1
    assert margin_slide[0]["status"] == "excluded"

    for s in outline["slides"]:
        if s["status"] == "active":
            s["narrative"] = {
                "observe": "Data shows 100 records.",
                "analyze": "Two regions present.",
                "synthesize": "Investigate Cost data quality.",
            }

    out = tmp_path / "messy.pptx"
    build_deck(outline, template_path=None, out_path=str(out))
    prs = Presentation(str(out))

    last = prs.slides[-1]
    text = "\n".join(s.text_frame.text for s in last.shapes if s.has_text_frame)
    assert "Exclusions" in text or "Integrity" in text
    assert "Cost" in text
