from pathlib import Path
from pptx import Presentation
from scripts.build_pptx import build_deck

def _outline_with_one_slide():
    return {
        "slides": [
            {
                "n": 1,
                "layout": "Title and Content",
                "title": "Executive Summary",
                "content_type": "exec_summary",
                "computation_id": None,
                "chart_spec": None,
                "status": "active",
                "narrative": {
                    "observe": "Q3 revenue is $4,200,000.",
                    "analyze": "Down 7% MoM.",
                    "synthesize": "Momentum stalled; investigate before Q4.",
                },
            }
        ]
    }

def test_build_deck_writes_pptx(tmp_path):
    outline = _outline_with_one_slide()
    out = tmp_path / "deck.pptx"
    result = build_deck(outline, template_path=None, out_path=str(out))
    assert Path(result["pptx_path"]).exists()
    prs = Presentation(result["pptx_path"])
    assert len(prs.slides) == 1

def test_build_deck_writes_synthesis_to_slide_body(tmp_path):
    outline = _outline_with_one_slide()
    out = tmp_path / "deck.pptx"
    build_deck(outline, template_path=None, out_path=str(out))
    prs = Presentation(str(out))
    slide = prs.slides[0]
    text = "\n".join(
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Momentum stalled" in text  # synthesize on slide body

def test_build_deck_writes_observe_analyze_to_speaker_notes(tmp_path):
    outline = _outline_with_one_slide()
    out = tmp_path / "deck.pptx"
    build_deck(outline, template_path=None, out_path=str(out))
    prs = Presentation(str(out))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "$4,200,000" in notes
    assert "7%" in notes

def test_build_deck_appends_exclusions_slide_when_skipped(tmp_path):
    outline = {
        "slides": [
            {
                "n": 1,
                "layout": "Title and Content",
                "title": "Active",
                "content_type": "exec_summary",
                "status": "active",
                "narrative": {"observe": "x is 1.", "analyze": "y.", "synthesize": "z."},
            },
            {
                "n": 2,
                "layout": "Image + Text",
                "title": "Margin",
                "content_type": "section",
                "status": "excluded",
                "reason": "Cost column 40% null",
            },
        ]
    }
    out = tmp_path / "deck.pptx"
    build_deck(outline, template_path=None, out_path=str(out))
    prs = Presentation(str(out))
    # 1 active + 1 exclusions slide
    assert len(prs.slides) == 2
    last = prs.slides[-1]
    text = "\n".join(s.text_frame.text for s in last.shapes if s.has_text_frame)
    assert "Exclusions" in text or "Integrity" in text
    assert "Margin" in text
    assert "Cost column 40% null" in text

def test_no_exclusions_slide_when_no_exclusions(tmp_path):
    outline = _outline_with_one_slide()
    out = tmp_path / "deck.pptx"
    build_deck(outline, template_path=None, out_path=str(out))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
