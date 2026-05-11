# presentation-builder/tests/test_e2e_pii.py
import re
from pathlib import Path
from pptx import Presentation
from scripts.ingest import ingest
from scripts.profile import profile, detect_context, build_outline
from scripts.build_pptx import build_deck

SSN_RX = re.compile(r"\d{3}-\d{2}-\d{4}")
EMAIL_RX = re.compile(r"[^@\s]+@[^@\s]+\.[^@\s]+")


def _all_text_in_deck(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide:
            parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


def test_pii_never_appears_in_deck(sample_pii, tmp_path):
    ing = ingest(str(sample_pii))
    df = ing["dataframe"]
    prof = profile(df)
    assert "SSN" in prof["pii_columns"]
    assert "Email" in prof["pii_columns"]

    ctx = detect_context(prof)
    outline = build_outline(prof, ctx)

    for slide in outline["slides"]:
        if slide["status"] == "active" and slide.get("computation_id"):
            slide["narrative"] = {
                "observe": "Revenue total available.",
                "analyze": "Across customers.",
                "synthesize": "Plan retention review.",
            }

    out = tmp_path / "deck_pii.pptx"
    build_deck(outline, template_path=None, out_path=str(out))

    text = _all_text_in_deck(str(out))
    assert not SSN_RX.search(text), "SSN leaked into deck"
    assert not EMAIL_RX.search(text), "Email leaked into deck"
    assert "Person 0" not in text and "Person 1" not in text
