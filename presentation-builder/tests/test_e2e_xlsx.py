# presentation-builder/tests/test_e2e_xlsx.py
from pathlib import Path
from pptx import Presentation
from scripts.ingest import ingest
from scripts.profile import profile
from scripts.context import detect_context
from scripts.outline import build_outline
from scripts.analyze import analyze
from scripts.aggregator import aggregate
from scripts.chart import render_chart
from scripts.narrative import validate_narrative
from scripts.build_pptx import build_deck


def test_e2e_xlsx_produces_complete_deck(sample_xlsx, tmp_path):
    # Stage 1
    ing = ingest(str(sample_xlsx))
    assert "dataframe" in ing
    df = ing["dataframe"]
    prof = profile(df)
    ctx = detect_context(prof)
    outline = build_outline(prof, ctx)
    assert outline["slides"][0]["content_type"] == "exec_summary"

    # Stage 2 (using a deterministic mock narrative — real LLM call would replace)
    for slide in outline["slides"]:
        if slide["status"] != "active":
            continue
        if slide.get("computation_id"):
            kv = analyze(df, slide["computation_id"])
            slide["kv"] = kv
            if kv:
                first_key = next(iter(kv))
                first_val = kv[first_key]
                # Only include numbers from kv in observe; keep analyze number-free
                narrative = {
                    "observe": f"{first_key} is {first_val}.",
                    "analyze": "Derived from available data points.",
                    "synthesize": "Review trend before next planning cycle.",
                }
                vr = validate_narrative(narrative, kv, pii_columns=prof.get("pii_columns"))
                assert vr["valid"], f"narrative invalid: {vr['mismatches']}"
                slide["narrative"] = narrative
        if slide["content_type"] == "section" and slide.get("computation_id") == "monthly_revenue":
            chart_data = aggregate(df, {"type": "line", "x": "Date", "y": "Revenue"})
            chart_path = tmp_path / f"chart_{slide['n']}.png"
            r = render_chart(chart_data, str(chart_path), title=slide["title"])
            if "png_path" in r:
                slide["chart_png"] = r["png_path"]

    # Exec summary narrative (static demo — not validated against KV)
    outline["slides"][0]["narrative"] = {
        "observe": "Pipeline processed 200 rows.",
        "analyze": "Across 4 regions and 3 products.",
        "synthesize": "Data is suitable for executive review.",
    }

    out = tmp_path / "deck.pptx"
    result = build_deck(outline, template_path=None, out_path=str(out))
    assert Path(result["pptx_path"]).exists()

    prs = Presentation(result["pptx_path"])
    assert len(prs.slides) >= 2
    title0 = prs.slides[0].shapes.title.text
    assert "Executive Summary" in title0 or "Summary" in title0
