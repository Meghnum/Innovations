from __future__ import annotations

from pathlib import Path
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError as _e:  # pragma: no cover
    raise ImportError(
        "presentation-builder requires 'python-pptx' — pip install python-pptx") from _e
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from render import pick_layout, get_brand_colors, get_brand_fonts  # sibling module; no package needed

DEFAULT_TEMPLATE = Path(__file__).parent.parent / "assets" / "default_template.pptx"

# small labels above the title, in the style of an exec deck "kicker"
_KICKER_FOR = {
    "exec_summary": "EXECUTIVE SUMMARY",
    "insight": "KEY INSIGHT",
    "section": "ANALYSIS",
    "deep_dive": "DEEP DIVE",
    "table": "DETAIL",
}
_INK = RGBColor(0x1F, 0x1F, 0x1F)   # near-black for titles/body text


def _resolve_template(template_path: str | None) -> str | None:
    """Explicit template if it exists, else the shipped default, else None
    (python-pptx then starts from its built-in blank template)."""
    if template_path and Path(template_path).exists():
        return template_path
    if DEFAULT_TEMPLATE.exists():
        return str(DEFAULT_TEMPLATE)
    return None


def _hex_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _style_for(template: str | None) -> dict:
    """Typography/colour tokens for the deck: theme fonts + brand colours from
    the template (or the documented fallbacks). One source for every slide —
    this is what keeps the deck's type consistent."""
    fonts = get_brand_fonts(template)
    colors = get_brand_colors(template)
    return {"display": fonts["display"], "body": fonts["body"],
            "primary": _hex_rgb(colors["primary"]), "ink": _INK}


def _write_para(p, text: str, font: str, size: float, bold: bool, color: RGBColor,
                space_after: int = 6):
    p.text = text
    p.space_after = Pt(space_after)
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_textbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(int(x), int(y), int(w), int(h)).text_frame
    tf.word_wrap = True
    return tf


_NOTES_BODY_SP = (
    '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<p:nvSpPr><p:cNvPr id="100" name="Notes Placeholder"/>'
    '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
    '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
    '<p:spPr><a:xfrm><a:off x="685800" y="1143000"/>'
    '<a:ext cx="5486400" cy="6400800"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
    '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr/></a:p></p:txBody></p:sp>'
)


def _notes_text_frame(notes):
    """The notes body text frame; some templates' notes masters carry no body
    placeholder, so cloned notes slides have none — inject one (speaker notes
    hold the provenance; they must never be silently dropped)."""
    if notes.notes_text_frame is not None:
        return notes.notes_text_frame
    from lxml import etree
    notes.shapes._spTree.append(etree.fromstring(_NOTES_BODY_SP))
    return notes.notes_text_frame


def _write_speaker_notes(slide, narrative: dict):
    parts = []
    if narrative.get("observe"):
        parts.append(f"Observe: {narrative['observe']}")
    if narrative.get("analyze"):
        parts.append(f"Analyze: {narrative['analyze']}")
    if narrative.get("synthesize"):
        parts.append(f"Synthesize: {narrative['synthesize']}")
    if not parts:
        return
    tf = _notes_text_frame(slide.notes_slide)
    if tf is not None:
        tf.text = "\n\n".join(parts)


def _body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            return shape
    return None


def _add_slide(prs, slide_data: dict, style: dict):
    """Render one slide on a consistent grid in the template's type:
    kicker · statement title · ≤3 body bullets (or the synthesize line) ·
    optional chart on the right half. Uses the template's own title/body
    placeholders when the layout provides them; otherwise draws styled
    textboxes (decks like ours have placeholder-less masters)."""
    layout = pick_layout(prs, slide_data.get("layout", "Title and Content"))
    slide = prs.slides.add_slide(layout)
    W, H = prs.slide_width, prs.slide_height
    margin = int(W * 0.055)
    content_w = W - 2 * margin

    narrative = slide_data.get("narrative", {})
    bullets = [str(b) for b in (slide_data.get("_bullets") or [])][:3]  # 3-bullet law
    if not bullets and narrative.get("synthesize"):
        bullets = [narrative["synthesize"]]

    # kicker — small bold label in the brand colour
    kicker = slide_data.get("kicker") or _KICKER_FOR.get(slide_data.get("content_type", ""), "")
    if slide_data.get("app"):
        kicker = f"{kicker}  ·  {slide_data['app']}" if kicker else str(slide_data["app"])
    if kicker:
        tf = _add_textbox(slide, margin, H * 0.055, content_w, H * 0.08)
        _write_para(tf.paragraphs[0], str(kicker).upper(), style["body"], 10, True,
                    style["primary"], space_after=0)

    # title — template placeholder if the layout has one, else display-font statement
    title = slide_data.get("title", "")
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    elif title:
        tf = _add_textbox(slide, margin, H * 0.13, content_w, H * 0.26)
        _write_para(tf.paragraphs[0], title, style["display"], 30, True, style["ink"],
                    space_after=0)

    # body — placeholder when available, else a styled textbox column
    has_chart = bool(slide_data.get("chart_png"))
    body = _body_placeholder(slide)
    if body and bullets:
        tf = body.text_frame
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _write_para(p, b, style["body"], 14, False, style["ink"])
    elif bullets:
        body_w = content_w * (0.50 if has_chart else 1.0)
        tf = _add_textbox(slide, margin, H * 0.42, body_w, H * 0.46)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            mark = "•  " if len(bullets) > 1 else ""
            _write_para(p, mark + b, style["body"], 13, False, style["ink"], space_after=10)

    if has_chart:
        slide.shapes.add_picture(
            slide_data["chart_png"],
            int(margin + content_w * 0.52), int(H * 0.30),
            width=int(content_w * 0.48),
        )

    if narrative:
        _write_speaker_notes(slide, narrative)


def _generate_exclusions_slide(prs, exclusions: list, style: dict):
    """Transparency slide: everything dropped for integrity/PII, in the same
    type as the rest of the deck (this list is not capped at 3 bullets —
    transparency beats brevity here)."""
    if not exclusions:
        return
    layout = pick_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    W, H = prs.slide_width, prs.slide_height
    margin = int(W * 0.055)
    content_w = W - 2 * margin

    tf = _add_textbox(slide, margin, H * 0.055, content_w, H * 0.08)
    _write_para(tf.paragraphs[0], "TRANSPARENCY", style["body"], 10, True,
                style["primary"], space_after=0)
    title = "Data Integrity & Exclusions Report"
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        tf = _add_textbox(slide, margin, H * 0.13, content_w, H * 0.20)
        _write_para(tf.paragraphs[0], title, style["display"], 28, True, style["ink"],
                    space_after=0)

    lines = ["The following sections were excluded or modified to ensure accuracy and compliance:"]
    lines += [f"•  {e['title']}: {e['reason']}" for e in exclusions]
    body = _body_placeholder(slide)
    tf = body.text_frame if body else _add_textbox(slide, margin, H * 0.38, content_w, H * 0.52)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _write_para(p, line, style["body"], 12, i == 0, style["ink"], space_after=8)


def build_deck(outline: dict, template_path: str | None, out_path: str) -> dict:
    """Assemble the deck. Returns {"pptx_path", "slide_count"}.
    All typography comes from one style block (template theme fonts/colours,
    or documented fallbacks). Excluded slides and any slide that fails to
    build are reported on the Exclusions slide instead of crashing the deck
    (SKILL.md failure mode: skip and log — never silent-fail)."""
    template = _resolve_template(template_path)
    prs = Presentation(template) if template else Presentation()
    style = _style_for(template)
    exclusions = []
    for slide_data in outline.get("slides", []):
        if slide_data.get("status") == "excluded":
            exclusions.append({
                "title": slide_data.get("title", "(untitled)"),
                "reason": slide_data.get("reason", "no reason provided"),
            })
            continue
        png = slide_data.get("chart_png")
        if png and not Path(png).exists():    # pre-check: keeps the deck free of half-built slides
            exclusions.append({
                "title": slide_data.get("title", "(untitled)"),
                "reason": f"chart image missing: {png}",
            })
            continue
        try:
            _add_slide(prs, slide_data, style)
        except Exception as e:
            exclusions.append({
                "title": slide_data.get("title", "(untitled)"),
                "reason": f"slide build failed: {e}",
            })
    _generate_exclusions_slide(prs, exclusions, style)
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return {"pptx_path": out_path, "slide_count": len(prs.slides)}
